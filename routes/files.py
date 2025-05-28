from flask import Blueprint, request, jsonify, session, current_app
from werkzeug.utils import secure_filename
import os
import uuid
import pdfplumber
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook
import filetype
from extensions import db, limiter
from models import File, User, ReadingLog
import re
import hashlib
import json
import traceback

files = Blueprint('files', __name__)

ALLOWED_EXTENSIONS = {'pdf', 'docx', 'pptx', 'txt', 'md', 'xlsx', 'xls', 'doc', 'ppt'}
ALLOWED_MIME_TYPES = {
    # PDF files
    'application/pdf': 'pdf',
    
    # Word documents
    'application/vnd.openxmlformats-officedocument.wordprocessingml.document': 'docx',
    'application/msword': 'doc',
    'application/zip': 'pptx',  # Modern Office files are detected as ZIP
    
    # PowerPoint files
    'application/vnd.openxmlformats-officedocument.presentationml.presentation': 'pptx',
    'application/vnd.ms-powerpoint': 'ppt',
    'application/powerpoint': 'ppt',
    'application/mspowerpoint': 'ppt',
    'application/x-mspowerpoint': 'ppt',
    
    # Excel files
    'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': 'xlsx',
    'application/vnd.ms-excel': 'xls',
    'application/excel': 'xls',
    'application/x-excel': 'xls',
    'application/x-msexcel': 'xls',
    
    # Text files
    'text/plain': 'txt',
    'text/markdown': 'md',
    'text/x-markdown': 'md'
}

def allowed_file(file):
    # First check file extension
    if '.' not in file.filename:
        print("No file extension found")
        return False
    
    extension = file.filename.rsplit('.', 1)[1].lower()
    if extension not in ALLOWED_EXTENSIONS:
        print(f"Extension {extension} not in allowed extensions: {ALLOWED_EXTENSIONS}")
        return False
    
    # Then verify actual file type using filetype
    file_head = file.read(4096)
    kind = filetype.guess(file_head)
    file.seek(0)  # Reset file pointer
    file_mime = kind.mime if kind else None
    print(f"File MIME type: {file_mime}")
    print(f"Allowed MIME types: {ALLOWED_MIME_TYPES}")
    
    # Special handling for Office files that are detected as ZIP
    if file_mime == 'application/zip' and extension in ['pptx', 'docx', 'xlsx']:
        print(f"Office file detected with extension {extension}")
        return True
    
    return file_mime in ALLOWED_MIME_TYPES

def extract_text(file_path, file_type):
    text = ""
    try:
        if file_type in ['pdf']:
            with pdfplumber.open(file_path) as pdf:
                text = "\n".join(page.extract_text() or "" for page in pdf.pages)
        
        elif file_type in ['docx', 'doc']:
            doc = Document(file_path)
            parts = []
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text.strip())
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = '\t'.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
            text = '\n'.join(parts)
        
        elif file_type in ['pptx', 'ppt']:
            prs = Presentation(file_path)
            slide_texts = []
            for idx, slide in enumerate(prs.slides, 1):
                slide_parts = [f"---- Slide {idx} ----"]
                for shape in slide.shapes:
                    # Extract text from shapes
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_parts.append(shape.text.strip())
                    # Extract text from tables
                    if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                        for row in shape.table.rows:
                            row_text = '\t'.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                            if row_text:
                                slide_parts.append(row_text)
                slide_texts.append('\n'.join(slide_parts))
            text = '\n\n'.join(slide_texts)
        
        elif file_type in ['xlsx', 'xls']:
            wb = load_workbook(file_path)
            for sheet in wb:
                text += f"\nSheet: {sheet.title}\n"
                for row in sheet.iter_rows():
                    row_text = " ".join(str(cell.value) for cell in row if cell.value)
                    if row_text:
                        text += row_text + "\n"
        
        elif file_type in ['txt', 'md']:
            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                text = f.read()
        
        else:
            print(f"Unsupported file type for text extraction: {file_type}")
            text = ""
        
        # Clean up special characters (replace non-printable with space)
        text = ''.join(c if c.isprintable() or c == '\n' or c == '\t' else ' ' for c in text)

    except Exception as e:
        print(f"Error extracting text from {file_path}: {str(e)}")
        text = ""
    
    # Calculate word count (split by whitespace and filter out empty strings)
    word_count = len([word for word in text.split() if word.strip()])
    
    return text, word_count

def calculate_hash(file_stream):
    hasher = hashlib.sha256()
    chunk_size = 4096
    while True:
        chunk = file_stream.read(chunk_size)
        if not chunk:
            break
        hasher.update(chunk)
    file_stream.seek(0)
    return hasher.hexdigest()

def process_tags(filename, text, provided_tags_json):
    if provided_tags_json:
        try:
            provided_tags = json.loads(provided_tags_json)
            if isinstance(provided_tags, list):
                return list(set(str(tag).strip() for tag in provided_tags if str(tag).strip()))
        except json.JSONDecodeError:
            pass
            
    filename_tags = set(re.findall(r'[A-Za-z]+', filename.lower()))
    return list(filename_tags)

@files.route('/upload', methods=['POST'])
@limiter.limit("20 per minute")
def upload_file():
    if 'user_id' not in session:
        return jsonify({'error': 'Not authenticated'}), 401
    
    user_id = session['user_id']
    print(f"Upload request from user_id: {user_id}")
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    print(f"Received file: {file.filename}")

    try:
        file_hash = calculate_hash(file.stream)
        print(f"File hash: {file_hash}")
    except Exception as e:
        print(f"Error calculating hash: {e}")
        return jsonify({'error': 'Could not calculate file hash'}), 500

    existing_file = File.query.filter_by(user_id=user_id, file_hash=file_hash).first()
    if existing_file:
        return jsonify({'error': f'Duplicate file detected. This content already exists as "{existing_file.name}".'}), 409

    # Check file type and get the normalized extension
    file_head = file.read(1024)
    kind = filetype.guess(file_head)
    file.seek(0)  # Reset file pointer
    file_mime = kind.mime if kind else None
    print(f"File MIME type: {file_mime}")
    # If file_mime is None, fall back to extension check
    if file_mime is None:
        extension = file.filename.rsplit('.', 1)[1].lower()
        if extension not in ALLOWED_EXTENSIONS:
            print(f"Unknown MIME type and extension {extension} not allowed.")
            return jsonify({'error': 'File type not allowed'}), 400
        normalized_extension = extension
    elif file_mime not in ALLOWED_MIME_TYPES:
        print(f"MIME type not allowed: {file_mime}")
        return jsonify({'error': 'File type not allowed'}), 400
    else:
        normalized_extension = ALLOWED_MIME_TYPES[file_mime]
    print(f"Normalized extension: {normalized_extension}")

    original_filename = secure_filename(file.filename)
    unique_filename = f"{str(uuid.uuid4())}.{normalized_extension}"
    
    # Create uploads directory in the working directory if it doesn't exist
    current_dir = os.getcwd()
    print(f"Current working directory: {current_dir}")
    
    upload_dir = os.path.join(current_dir, 'uploads')
    user_upload_dir = os.path.join(upload_dir, str(user_id))
    print(f"Creating user upload directory: {user_upload_dir}")
    
    os.makedirs(user_upload_dir, exist_ok=True)
    
    # Save file with relative path
    relative_file_path = os.path.join('uploads', str(user_id), unique_filename)
    absolute_file_path = os.path.join(current_dir, relative_file_path)
    print(f"Saving file to: {absolute_file_path}")
    
    try:
        file.save(absolute_file_path)
        print(f"File saved successfully")
    except Exception as e:
        print(f"Error saving file: {e}")
        return jsonify({'error': 'Failed to save file on server'}), 500

    # Verify file was saved
    if os.path.exists(absolute_file_path):
        print(f"File exists at {absolute_file_path}")
        print(f"File size: {os.path.getsize(absolute_file_path)} bytes")
    else:
        print(f"File does not exist at {absolute_file_path}")
        return jsonify({'error': 'File was not saved properly'}), 500

    extracted_text = ""
    word_count = 0
    try:
        extracted_text, word_count = extract_text(absolute_file_path, normalized_extension)
        print(f"Extracted {word_count} words from file")
    except Exception as e:
        print(f"Error extracting text from {original_filename}: {e}")

    provided_tags_json = request.form.get('tags')
    tags = process_tags(original_filename, extracted_text, provided_tags_json)
    
    # Calculate estimated duration based on word count
    # Assuming average reading speed of 200 words per minute
    estimated_duration = (word_count / 200) * 60 if word_count > 0 else 3600
    
    # Store the relative path in the database
    db_file = File(
        name=original_filename,
        filetype=normalized_extension,
        path=relative_file_path,  # Store relative path instead of absolute
        file_hash=file_hash,
        extracted_text=extracted_text,
        word_count=word_count,
        tags=tags,
        user_id=user_id,
        estimated_duration=int(estimated_duration)
    )
    
    try:
        db.session.add(db_file)
        db.session.commit()
        print(f"File record saved to database with ID: {db_file.id}")
    except Exception as e:
        db.session.rollback()
        print(f"Database error during file upload: {e}")
        try:
            os.remove(absolute_file_path)
            print(f"Cleaned up file after database error")
        except OSError as remove_error:
            print(f"Error cleaning up file {absolute_file_path} after DB error: {remove_error}")
        return jsonify({'error': 'Database error saving file record'}), 500

    return jsonify({
        'message': 'File uploaded successfully',
        'file': {
            'id': db_file.id,
            'name': db_file.name,
            'type': db_file.filetype,
            'tags': db_file.tags,
            'hash': db_file.file_hash,
            'word_count': db_file.word_count,
            'estimated_duration': db_file.estimated_duration
        }
    })

@files.route('/check-file', methods=['POST'])
def check_file_exists():
    if 'user_id' not in session:
        return jsonify({'error': 'Not authenticated'}), 401
    
    if 'file' not in request.files:
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400

    try:
        file_hash = calculate_hash(file.stream)
        existing_file = File.query.filter_by(
            user_id=session['user_id'], 
            file_hash=file_hash
        ).first()
        
        if existing_file:
            return jsonify({
                'exists': True,
                'file': {
                    'id': existing_file.id,
                    'name': existing_file.name,
                    'type': existing_file.filetype,
                    'upload_date': existing_file.upload_date.isoformat()
                }
            })
        
        return jsonify({'exists': False})
    
    except Exception as e:
        print(f"Error checking file: {e}")
        return jsonify({'error': 'Could not check file'}), 500

@files.route('/files', methods=['GET'])
def get_files():
    if 'user_id' not in session:
        return jsonify({'error': 'Not authenticated'}), 401
    
    files = File.query.filter_by(user_id=session['user_id']).order_by(File.upload_date.desc()).all()
    return jsonify({
        'files': [{
            'id': file.id,
            'name': file.name,
            'filetype': file.filetype,
            'tags': file.tags,
            'upload_date': file.upload_date.isoformat(),
            'estimated_duration': file.estimated_duration,
            'path': file.path if file.path else None
        } for file in files]
    })

@files.route('/files/<int:file_id>', methods=['GET'])
def get_file(file_id):
    if 'user_id' not in session:
        return jsonify({'error': 'Not authenticated'}), 401
    
    file = File.query.filter_by(id=file_id, user_id=session['user_id']).first_or_404()
    
    return jsonify({
        'file': {
            'id': file.id,
            'name': file.name,
            'filetype': file.filetype,
            'tags': file.tags,
            'upload_date': file.upload_date.isoformat(),
            'text_content': file.extracted_text,
            'hash': file.file_hash
        }
    })

@files.route('/update-word-counts', methods=['POST'])
def update_word_counts():
    if 'user_id' not in session:
        return jsonify({'error': 'Not authenticated'}), 401
    
    try:
        # Get all files for the current user
        files = File.query.filter_by(user_id=session['user_id']).all()
        updated_count = 0
        
        for file in files:
            if file.path and os.path.exists(file.path):
                try:
                    # Extract text and get word count
                    extracted_text, word_count = extract_text(file.path, file.filetype)
                    
                    # Update file record
                    file.word_count = word_count
                    file.estimated_duration = int((word_count / 200) * 60) if word_count > 0 else 3600
                    updated_count += 1
                except Exception as e:
                    print(f"Error processing file {file.name}: {e}")
                    continue
        
        db.session.commit()
        return jsonify({
            'message': f'Successfully updated {updated_count} files',
            'updated_count': updated_count
        })
    
    except Exception as e:
        db.session.rollback()
        print(f"Error updating word counts: {e}")
        return jsonify({'error': 'Failed to update word counts'}), 500

@files.route('/<int:file_id>', methods=['DELETE'])
def delete_file_route(file_id):
    print(f"=== DELETE FILE API CALLED for file_id: {file_id} ===")
    if 'user_id' not in session:
        print(f"Error: User not authenticated")
        return jsonify({'error': 'Not authenticated'}), 401

    user_id = session['user_id']
    print(f"User ID from session: {user_id}")

    # Find the file belonging to the user
    file_to_delete = File.query.filter_by(id=file_id, user_id=user_id).first()

    if not file_to_delete:
        print(f"Error: File {file_id} not found for user {user_id}")
        return jsonify({'error': 'File not found or you do not have permission to delete it'}), 404

    file_path = file_to_delete.path
    file_name = file_to_delete.name
    print(f"Found file to delete: ID={file_to_delete.id}, Name='{file_name}', Path='{file_path}'")

    try:
        # 1. Delete associated ReadingLog entries
        logs_deleted_count = ReadingLog.query.filter_by(file_id=file_id, user_id=user_id).delete()
        print(f"Deleted {logs_deleted_count} associated ReadingLog entries.")

        # 2. Delete the physical file (if path exists)
        if file_path and os.path.exists(file_path):
            try:
                os.remove(file_path)
                print(f"Successfully deleted physical file: {file_path}")
            except OSError as e:
                print(f"Error deleting physical file {file_path}: {e}")
                # Log this error but proceed to delete DB record maybe? Or return error?
                # For now, let's return an error if physical file deletion fails.
                return jsonify({'error': f'Failed to delete physical file: {e.strerror}'}), 500
        elif file_path:
            print(f"Warning: Physical file path recorded but not found: {file_path}")

        # 3. Delete the File record from the database
        db.session.delete(file_to_delete)
        print(f"Marked file record ID {file_id} for deletion.")

        # 4. Commit all changes
        db.session.commit()
        print(f"Successfully committed deletion for file ID {file_id}.")

        return jsonify({'message': f"Successfully deleted file '{file_name}'"}), 200

    except Exception as e:
        db.session.rollback()
        print(f"Error during file deletion process for file ID {file_id}: {e}")
        traceback.print_exc()
        return jsonify({'error': f'An unexpected error occurred during file deletion: {str(e)}'}), 500 